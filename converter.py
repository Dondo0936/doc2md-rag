# converter.py — Document → Markdown conversion with post-processing
#
# Pipeline: Extract raw markdown → Convert tables to lists → Handle image descriptions
# Backends: docling (primary) → pypandoc/pymupdf4llm/markitdown (fallback)

import logging
import re
import os
import tempfile
import pandas as pd
from html.parser import HTMLParser

logger = logging.getLogger(__name__)

# ─── Docling (Primary Backend) ──────────────────────────────────

# Lazy singletons keyed by (ocr, describe_images)
_docling_converters: dict = {}


def _get_docling_converter(ocr: bool = False, describe_images: bool = False):
    """Lazy-init a Docling DocumentConverter singleton.

    Args:
        ocr: If True, enable RapidOCR for scanned PDFs.
        describe_images: If True, enable SmolVLM image descriptions.
    """
    global _docling_converters
    key = (ocr, describe_images)
    if key in _docling_converters:
        return _docling_converters[key]

    from docling.document_converter import (
        DocumentConverter, PdfFormatOption, WordFormatOption, PowerpointFormatOption,
    )
    from docling.datamodel.base_models import InputFormat
    from docling.datamodel.pipeline_options import (
        PdfPipelineOptions,
        ConvertPipelineOptions,
        TableStructureOptions,
        TableFormerMode,
    )

    # --- PDF pipeline ---
    pdf_options = PdfPipelineOptions()
    pdf_options.do_table_structure = True
    pdf_options.table_structure_options = TableStructureOptions(
        mode=TableFormerMode.ACCURATE,
        do_cell_matching=True,
    )

    if describe_images:
        pdf_options.generate_picture_images = True
        pdf_options.do_picture_description = True
        logger.info("PDF image descriptions enabled via SmolVLM")
    else:
        pdf_options.generate_picture_images = False

    if ocr:
        pdf_options.do_ocr = True
        try:
            from docling.datamodel.pipeline_options import RapidOcrOptions
            pdf_options.ocr_options = RapidOcrOptions()
            logger.info("RapidOCR enabled for scanned PDF processing")
        except ImportError:
            logger.warning("RapidOcrOptions not available; using Docling default OCR")
    else:
        pdf_options.do_ocr = False

    format_options = {
        InputFormat.PDF: PdfFormatOption(pipeline_options=pdf_options),
    }

    # --- DOCX / PPTX pipeline (image descriptions via SmolVLM) ---
    if describe_images:
        doc_options = ConvertPipelineOptions(do_picture_description=True)
        format_options[InputFormat.DOCX] = WordFormatOption(pipeline_options=doc_options)
        format_options[InputFormat.PPTX] = PowerpointFormatOption(pipeline_options=doc_options)
        logger.info("DOCX/PPTX image descriptions enabled via SmolVLM")

    converter = DocumentConverter(format_options=format_options)
    _docling_converters[key] = converter
    return converter


def _extract_with_docling(file_path: str, ocr: bool = False, describe_images: bool = False) -> str:
    """Extract any supported format using Docling. Returns markdown string.

    Args:
        ocr: If True, use the OCR-enabled converter (for scanned PDFs).
        describe_images: If True, use the image-description converter.
    """
    converter = _get_docling_converter(ocr=ocr, describe_images=describe_images)
    result = converter.convert(file_path, raises_on_error=False)

    from docling.datamodel.document import ConversionStatus
    if result.status == ConversionStatus.FAILURE:
        errors = "; ".join(str(e) for e in (result.errors or []))
        raise RuntimeError(f"Docling conversion failed: {errors}")

    md = result.document.export_to_markdown()
    if not md or len(md.strip()) < 10:
        raise RuntimeError("Docling produced empty/minimal output")

    # When image descriptions are enabled, Docling still emits residual
    # <!-- image --> placeholders after the description text — strip them.
    if describe_images:
        md = re.sub(r"\n*<!-- image -->\n*", "\n\n", md)

    return md


def warmup_converter():
    """Pre-initialize the default (non-OCR) Docling converter.

    Call this at app startup to avoid cold-start delay on first upload.
    Returns the converter instance (useful for @st.cache_resource).
    """
    return _get_docling_converter(ocr=False)


MAX_FILE_SIZE_MB = 100
MAX_FILE_SIZE_BYTES = MAX_FILE_SIZE_MB * 1024 * 1024


# ─── HTML Table → Pipe Table Converter ────────────────────────

class _HTMLTableParser(HTMLParser):
    """Parse HTML tables into a list of rows (list of cell strings)."""
    def __init__(self):
        super().__init__()
        self.tables = []
        self._table = None
        self._row = None
        self._cell = None
        self._in_cell = False

    def handle_starttag(self, tag, attrs):
        if tag == "table":
            self._table = []
        elif tag == "tr":
            self._row = []
        elif tag in ("td", "th"):
            self._cell = []
            self._in_cell = True
        elif tag == "p" and self._in_cell and self._cell:
            self._cell.append("<br>")

    def handle_endtag(self, tag):
        if tag == "table" and self._table is not None:
            self.tables.append(self._table)
            self._table = None
        elif tag == "tr" and self._row is not None:
            if self._table is not None:
                self._table.append(self._row)
            self._row = None
        elif tag in ("td", "th") and self._cell is not None:
            text = "".join(self._cell).strip()
            text = re.sub(r"<br>\s*$", "", text)
            if self._row is not None:
                self._row.append(text)
            self._cell = None
            self._in_cell = False

    def handle_data(self, data):
        if self._in_cell and self._cell is not None:
            self._cell.append(data)


def _html_tables_to_pipe(md_text: str) -> str:
    """Convert HTML <table> blocks in markdown to pipe-delimited tables."""
    table_re = re.compile(r"<table[^>]*>.*?</table>", re.DOTALL | re.IGNORECASE)

    def _convert(match):
        parser = _HTMLTableParser()
        parser.feed(match.group(0))
        if not parser.tables:
            return match.group(0)
        table = parser.tables[0]
        if not table:
            return ""
        max_cols = max(len(r) for r in table)
        lines = []
        for i, row in enumerate(table):
            while len(row) < max_cols:
                row.append("")
            cells = [c.replace("|", "∣").replace("\n", "<br>") for c in row]
            lines.append("| " + " | ".join(cells) + " |")
            if i == 0:
                lines.append("| " + " | ".join(["---"] * max_cols) + " |")
        return "\n".join(lines)

    return table_re.sub(_convert, md_text)


# ─── Password / Scanned PDF Detection ────────────────────────

def check_pdf_encrypted(file_path: str) -> dict:
    """Check if a PDF is password-protected.

    Returns:
        {"encrypted": False} — not encrypted or only owner-password (can open freely)
        {"encrypted": True}  — needs a user password to open
    """
    try:
        import pikepdf
    except ImportError:
        return {"encrypted": False}

    try:
        pdf = pikepdf.open(file_path, password="")
        pdf.close()
        return {"encrypted": False}
    except pikepdf.PasswordError:
        return {"encrypted": True}
    except Exception:
        return {"encrypted": False}


def decrypt_pdf(file_path: str, password: str) -> str:
    """Decrypt a password-protected PDF with the given password.

    Returns path to a decrypted temporary file.
    Raises ValueError if the password is wrong.
    """
    import pikepdf

    try:
        pdf = pikepdf.open(file_path, password=password)
    except pikepdf.PasswordError:
        raise ValueError("Incorrect password")

    decrypted = file_path + ".decrypted.pdf"
    pdf.save(decrypted)
    pdf.close()
    return decrypted


def _try_decrypt_pdf(file_path: str) -> str:
    """Decrypt owner-protected PDFs (no user password needed)."""
    try:
        import pikepdf
    except ImportError:
        return file_path
    try:
        pdf = pikepdf.open(file_path, password="")
        decrypted = file_path + ".decrypted.pdf"
        pdf.save(decrypted)
        pdf.close()
        return decrypted
    except Exception:
        return file_path


def _is_scanned_pdf(file_path: str, sample_pages: int = 3) -> bool:
    """Heuristic: pages with <30 chars text but images → scanned."""
    import fitz
    doc = fitz.open(file_path)
    pages = min(len(doc), sample_pages)
    scanned = 0
    for i in range(pages):
        page = doc[i]
        text = page.get_text("text").strip()
        images = page.get_images(full=True)
        if len(text) < 30 and len(images) > 0:
            scanned += 1
    doc.close()
    return scanned >= max(1, pages * 0.6)


def _fitz_table_for_page(page) -> str:
    """Fallback: extract tables from a fitz page using find_tables()."""
    tables = page.find_tables()
    if not tables.tables:
        return ""
    main_table = max(tables.tables, key=lambda t: len(t.cells))
    data = main_table.extract()
    if not data:
        return ""
    active_cols = sorted({
        j for row in data for j, v in enumerate(row) if v and v.strip()
    })
    if not active_cols:
        return ""
    clean_rows = []
    for row in data:
        clean_row = []
        for j in active_cols:
            val = row[j] if j < len(row) and row[j] else ""
            val = val.replace("\n", "<br>").replace("|", "∣").strip()
            clean_row.append(val)
        clean_rows.append(clean_row)
    num_cols = len(active_cols)
    lines = ["| " + " | ".join(clean_rows[0]) + " |"]
    lines.append("| " + " | ".join(["---"] * num_cols) + " |")
    for row in clean_rows[1:]:
        while len(row) < num_cols:
            row.append("")
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines) + "\n"


def _extract_pdf(file_path: str, embed_images: bool = False) -> str:
    import pymupdf4llm
    import fitz

    usable = _try_decrypt_pdf(file_path)
    try:
        if _is_scanned_pdf(usable):
            logger.warning("Scanned PDF detected — text extraction may be incomplete.")

        md_pages = pymupdf4llm.to_markdown(usable, page_chunks=True,
                                            embed_images=embed_images)
        doc = fitz.open(usable)
        parts = []
        for idx, chunk in enumerate(md_pages):
            page_text = chunk["text"].strip()
            meta = chunk.get("metadata", {})
            page_num = meta.get("page", meta.get("page_number", idx + 1))
            page_idx = page_num - 1 if page_num >= 1 else page_num
            if len(page_text) <= 10 and 0 <= page_idx < len(doc):
                table_md = _fitz_table_for_page(doc[page_idx])
                if table_md:
                    parts.append(table_md)
                    continue
            parts.append(chunk["text"])
        doc.close()
        return "\n\n".join(parts)
    finally:
        if usable != file_path and os.path.exists(usable):
            os.unlink(usable)


def _extract_docx(file_path: str) -> str:
    """Extract DOCX using pypandoc (handles math, track changes, footnotes).
    Falls back to python-docx if pypandoc is unavailable."""
    try:
        import pypandoc
        md = pypandoc.convert_file(
            file_path, "gfm",
            extra_args=["--track-changes=accept", "--wrap=none",
                        "--markdown-headings=atx"],
        )
        md = _html_tables_to_pipe(md)
        return md
    except Exception as e:
        logger.warning("pypandoc failed (%s), falling back to python-docx", e)
        return _extract_docx_legacy(file_path)


def _extract_docx_legacy(file_path: str) -> str:
    """Legacy DOCX extraction using python-docx (fallback)."""
    from docx import Document
    from docx.table import Table
    from docx.text.paragraph import Paragraph
    from docx.oxml.ns import qn

    doc = Document(file_path)
    parts = []

    def _process_paragraph(para):
        style = para.style.name if para.style else ""
        text = para.text.strip()
        if not text:
            parts.append("")
            return
        if style.startswith("Heading 1"):
            parts.append(f"# {text}")
        elif style.startswith("Heading 2"):
            parts.append(f"## {text}")
        elif style.startswith("Heading 3"):
            parts.append(f"### {text}")
        elif style.startswith("List"):
            parts.append(f"- {text}")
        else:
            m = re.match(r"^(\d+(?:\.\d+){2,})\s*\.?\s+(.+)", text)
            if m:
                depth = m.group(1).count(".")
                level = min(depth + 1, 6)
                parts.append(f"{'#' * level} {text}")
            else:
                parts.append(text)

    def _sanitize_cell(text):
        text = text.strip()
        text = text.replace("\n", "<br>")
        text = text.replace("|", "∣")
        return text

    def _process_table(table):
        if not table.rows:
            return
        def _is_merged_row(row):
            texts = [cell.text.strip() for cell in row.cells]
            unique = set(t for t in texts if t)
            return len(unique) == 1 and unique
        header_idx = 0
        for idx, row in enumerate(table.rows):
            if _is_merged_row(row):
                text = next(t for t in (cell.text.strip() for cell in row.cells) if t)
                parts.append("")
                parts.append(f"### {text}")
                header_idx = idx + 1
            else:
                break
        if header_idx >= len(table.rows):
            return
        headers = [_sanitize_cell(cell.text) for cell in table.rows[header_idx].cells]
        parts.append("")
        parts.append("| " + " | ".join(headers) + " |")
        parts.append("| " + " | ".join(["---"] * len(headers)) + " |")
        for row in table.rows[header_idx + 1:]:
            if _is_merged_row(row):
                text = next(t for t in (cell.text.strip() for cell in row.cells) if t)
                parts.append("")
                parts.append(f"### {text}")
                parts.append("")
                parts.append("| " + " | ".join(headers) + " |")
                parts.append("| " + " | ".join(["---"] * len(headers)) + " |")
                continue
            cells = [_sanitize_cell(cell.text) for cell in row.cells]
            parts.append("| " + " | ".join(cells) + " |")
        parts.append("")

    for element in doc.element.body:
        if element.tag == qn("w:p"):
            _process_paragraph(Paragraph(element, doc))
        elif element.tag == qn("w:tbl"):
            _process_table(Table(element, doc))

    return "\n".join(parts)


def _extract_pptx(file_path: str) -> str:
    """Extract PPTX using MarkItDown (speaker notes, charts, grouped shapes).
    Falls back to python-pptx if MarkItDown is unavailable."""
    try:
        from markitdown import MarkItDown
        mid = MarkItDown()
        result = mid.convert(file_path)
        return result.text_content
    except Exception as e:
        logger.warning("MarkItDown failed (%s), falling back to python-pptx", e)
        return _extract_pptx_legacy(file_path)


def _extract_pptx_legacy(file_path: str) -> str:
    """Legacy PPTX extraction using python-pptx (fallback)."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    prs = Presentation(file_path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"## Slide {i}")
        for shape in slide.shapes:
            # Handle embedded images (guard against linked images without blob)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    _ = shape.image  # test access — linked images raise ValueError
                    name = shape.name or "visual content"
                    parts.append(f"[Image: {name}]")
                except (ValueError, AttributeError):
                    parts.append("[Image: linked image not available]")
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    # \x0b = vertical tab = soft line break in PPTX
                    text = para.text.replace("\x0b", "\n").strip()
                    if text:
                        parts.append(text)
            if shape.has_table:
                table = shape.table

                def _pptx_sanitize(text):
                    return text.strip().replace("\n", "<br>").replace("|", "∣")

                headers = [_pptx_sanitize(cell.text) for cell in table.rows[0].cells]
                parts.append("")
                parts.append("| " + " | ".join(headers) + " |")
                parts.append("| " + " | ".join(["---"] * len(headers)) + " |")
                for row in list(table.rows)[1:]:
                    cells = [_pptx_sanitize(cell.text) for cell in row.cells]
                    parts.append("| " + " | ".join(cells) + " |")
                parts.append("")
        parts.append("")
    return "\n".join(parts)


def _extract_csv(file_path: str) -> str:
    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            df = pd.read_csv(file_path, encoding=encoding)
            return df.to_markdown(index=False)
        except UnicodeDecodeError:
            continue
    # Last resort: replace errors
    df = pd.read_csv(file_path, encoding="utf-8", errors="replace")
    return df.to_markdown(index=False)


def _collapse_spaced_caps(text: str) -> str:
    """Collapse spaced-out capital headers: 'C A P I T A L' → 'CAPITAL'.

    Also collapses multi-space gaps left after merging adjacent words.
    """
    def _collapse(m):
        return m.group(0).replace(" ", "")
    # Match 3+ single uppercase letters separated by single spaces
    text = re.sub(r"\b[A-Z](?:\s[A-Z]){2,}\b", _collapse, text)
    # Collapse runs of 2+ spaces (left between collapsed words) to single space
    text = re.sub(r"  +", " ", text)
    return text


def _strip_page_numbers(text: str) -> str:
    """Remove standalone page numbers (1-4 digit number alone between blank lines)."""
    return re.sub(r"\n\n\s*\d{1,4}\s*\n\n", "\n\n", text)


def _extract_markdown(file_path: str, llm_client=None) -> str:
    """Extract markdown from PDF, DOCX, PPTX, or CSV.

    Uses Docling as primary backend for PDF/DOCX/PPTX, with legacy fallbacks.
    """
    file_size = os.path.getsize(file_path)
    if file_size > MAX_FILE_SIZE_BYTES:
        raise ValueError(f"File too large ({file_size / 1024 / 1024:.1f}MB). Maximum is {MAX_FILE_SIZE_MB}MB.")

    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".csv":
        return _extract_csv(file_path)
    if ext not in (".pdf", ".docx", ".pptx"):
        raise ValueError(f"Unsupported file type: {ext}")

    # Detect scanned PDFs for OCR routing
    use_ocr = False
    if ext == ".pdf":
        try:
            use_ocr = _is_scanned_pdf(file_path)
            if use_ocr:
                logger.info("Scanned PDF detected — enabling OCR via RapidOCR")
        except Exception as e:
            logger.debug("Scanned PDF check failed: %s", e)

    # Enable image descriptions for PDFs (SmolVLM — free, local)
    # Image descriptions via SmolVLM — enabled when ENABLE_IMAGE_DESC=1.
    # Needs ~1-2GB RAM; works on HF Spaces (16GB) but not Streamlit Cloud (1GB).
    describe_imgs = (
        os.environ.get("ENABLE_IMAGE_DESC", "0") == "1"
        and ext in (".pdf", ".docx", ".pptx")
    )

    # Try Docling first for PDF, DOCX, PPTX
    try:
        md = _extract_with_docling(file_path, ocr=use_ocr, describe_images=describe_imgs)
        md = _collapse_spaced_caps(md)
        md = _strip_page_numbers(md)
        logger.info("Docling extraction succeeded for %s%s%s", ext,
                     " (with OCR)" if use_ocr else "",
                     " (with image descriptions)" if describe_imgs else "")
        return md
    except Exception as e:
        logger.warning("Docling failed (%s), falling back to legacy backend", e)

    # Legacy fallbacks
    if ext == ".pdf":
        md = _extract_pdf(file_path, embed_images=bool(llm_client))
    elif ext == ".docx":
        md = _extract_docx(file_path)
    elif ext == ".pptx":
        md = _extract_pptx(file_path)

    md = _collapse_spaced_caps(md)
    md = _strip_page_numbers(md)
    return md


def _clean_cell(text: str) -> str:
    """Clean a table cell: convert <br> to newlines, strip markdown artifacts.

    Rejoins word-wrap fragments from narrow PDF columns while preserving
    genuine line breaks (bullets, sentence boundaries).
    """
    text = re.sub(r"\u200b", "", text)        # zero-width space
    text = text.strip()
    text = re.sub(r"<br\s*/?>", "\n", text)   # <br> → newline
    # Convert bullet markers to plain dashes
    text = re.sub(r"^[●○]\s*", "- ", text, flags=re.MULTILINE)
    # Rejoin lines that are just word-wrap artifacts from narrow PDF columns
    lines = text.split("\n")
    if len(lines) > 1:
        merged = [lines[0]]
        in_list = False  # after a colon-ended line, keep items separate
        for line in lines[1:]:
            stripped = line.strip()
            if not stripped:
                merged.append("")
                in_list = False
                continue
            # Keep as new line if it starts with a bullet/list marker
            if re.match(r"^[-●○*■□▪▸►•]\s", stripped) or re.match(r"^\d+[.)]\s", stripped):
                merged.append(stripped)
                in_list = False
                continue
            prev = merged[-1].rstrip() if merged else ""
            # After a colon-ended line, keep subsequent items on separate lines
            if in_list:
                merged.append(stripped)
                if stripped.endswith((".", ":", "!", "?", ";")):
                    in_list = stripped.endswith(":")
                continue
            # Join if previous line doesn't end with sentence punctuation
            if prev and not prev.endswith((".", ":", "!", "?", ";")) and not prev.endswith("  "):
                # Detect mid-word/identifier breaks: both fragments have no spaces
                # and at least one contains identifier chars like _ or .
                prev_token = prev.split("\n")[-1].strip()
                prev_is_token = " " not in prev_token
                cur_is_token = " " not in stripped
                has_id_char = ("_" in prev_token or "_" in stripped
                               or "." in prev_token or "." in stripped)
                if (prev_is_token and cur_is_token and has_id_char
                        and prev[-1].isalnum() and stripped[0].isalnum()):
                    merged[-1] = prev + stripped
                else:
                    merged[-1] = prev + " " + stripped
            else:
                merged.append(stripped)
                if prev.endswith(":"):
                    in_list = True
        text = "\n".join(merged)
    return text


def _parse_pipe_row(line: str) -> list[str]:
    """Parse a single pipe-delimited row into cells, handling <br> within cells."""
    line = line.strip()
    if line.startswith("|"):
        line = line[1:]
    if line.endswith("|"):
        line = line[:-1]
    return [_clean_cell(c) for c in line.split("|")]


def _is_separator(line: str) -> bool:
    """Check if a line is a markdown table separator (|---|---|).

    Each cell must be only dashes with optional leading/trailing colons
    (for alignment), e.g. |---|, |:---|, |:---:|, |---:|.
    This avoids false-positives on data rows containing dashes or spaces.
    """
    stripped = line.strip()
    if not stripped.startswith("|") or not stripped.endswith("|"):
        return False
    # Remove outer pipes, split into cells
    inner = stripped[1:-1]
    if not inner:
        return False
    cells = inner.split("|")
    # Every cell must match the separator pattern: optional colons around 3+ dashes
    # (markdown spec requires at least 3 dashes per cell)
    return all(re.match(r"^\s*:?-{3,}:?\s*$", cell) for cell in cells)


def _table_to_list(lines: list[str]) -> str:
    """Convert parsed table lines into a nested bullet list."""
    if not lines:
        return ""

    # Strip leading whitespace from each line (handle indented tables)
    lines = [l.strip() for l in lines]

    # Find separator line
    sep_idx = -1
    for i, line in enumerate(lines):
        if _is_separator(line):
            sep_idx = i
            break

    if sep_idx >= 0:
        # Header row is above separator
        headers = _parse_pipe_row(lines[sep_idx - 1]) if sep_idx > 0 else []
        # The row above separator might also be a data row if headers are generic
        # (e.g. "Col1", "Col2") — detect and include it as data
        generic_headers = all(
            re.match(r"^(Col\s*\d+|)$", h, re.IGNORECASE)
            for h in headers
        )

        data_start = sep_idx + 1
        # Collect any rows above the header as data
        pre_header_rows = []
        if sep_idx > 1:
            for i in range(0, sep_idx - 1):
                if lines[i].startswith("|") and not _is_separator(lines[i]):
                    pre_header_rows.append(lines[i])
    else:
        # No separator — treat first row as header
        headers = _parse_pipe_row(lines[0])
        generic_headers = True
        data_start = 1
        pre_header_rows = []

    # Collect data rows (skip any additional separators)
    data_rows = []
    for row_line in pre_header_rows:
        data_rows.append(_parse_pipe_row(row_line))

    if generic_headers and sep_idx > 0:
        # The "header" is actually data — include it
        data_rows.append(_parse_pipe_row(lines[sep_idx - 1]))

    for line in lines[data_start:]:
        if _is_separator(line) or not line.strip():
            continue
        data_rows.append(_parse_pipe_row(line))

    if not data_rows:
        return ""

    # Determine real column names
    if generic_headers or not headers:
        # Use Col1, Col2, ... as fallback
        max_cols = max(len(r) for r in data_rows) if data_rows else 1
        headers = [f"Col {i+1}" for i in range(max_cols)]

    # Merge continuation rows: if a row has empty first cell, append its
    # content to the previous row's cells (common in PDF tables with merged cells)
    merged_rows = []
    for row in data_rows:
        first_cell = row[0].strip() if row else ""
        if not first_cell and merged_rows:
            # Continuation of previous row — merge cell contents
            prev = merged_rows[-1]
            for j in range(len(row)):
                val = row[j].strip() if j < len(row) else ""
                if not val:
                    continue
                if j < len(prev) and prev[j].strip():
                    prev[j] = prev[j].rstrip() + "\n" + val
                elif j < len(prev):
                    prev[j] = val
                else:
                    prev.extend([""] * (j - len(prev) + 1))
                    prev[j] = val
        else:
            merged_rows.append(list(row))
    data_rows = merged_rows

    result_lines = []
    for i, row in enumerate(data_rows):
        # Detect merged rows: all non-empty cells identical → section heading
        non_empty = [c.strip() for c in row if c.strip()]
        unique_vals = set(non_empty)
        if len(unique_vals) == 1 and len(non_empty) > 1:
            result_lines.append(f"\n### {non_empty[0]}\n")
            continue

        # Use first cell as label; if empty, try second cell; last resort "Row N"
        label = row[0].split("\n")[0].strip() if row and row[0].strip() else ""
        if not label and len(row) > 1 and row[1].strip():
            label = row[1].split("\n")[0].strip()[:80]
        if not label:
            label = f"Row {i+1}"
        result_lines.append(f"- **{label}**:")
        for j, header in enumerate(headers):
            value = row[j] if j < len(row) else ""
            if not value or not header.strip():
                continue
            # Skip only if value is identical to label (avoid repeating the row header)
            if j == 0 and value.split("\n")[0].strip() == label:
                # Still emit remaining lines if multi-line
                extra_lines = value.split("\n")[1:]
                extra_lines = [vl.strip() for vl in extra_lines if vl.strip()]
                if extra_lines:
                    result_lines.append(f"  - {header}:")
                    for vl in extra_lines:
                        result_lines.append(f"    - {vl}")
                continue
            # If label came from this column (j==1 fallback), skip repeating it
            if j == 1 and not row[0].strip() and value.split("\n")[0].strip()[:80] == label:
                extra_lines = [vl.strip() for vl in value.split("\n")[1:] if vl.strip()]
                if extra_lines:
                    result_lines.append(f"  - {header}:")
                    for vl in extra_lines:
                        result_lines.append(f"    - {vl}")
                continue
            # Multi-line values get indented
            value_lines = value.split("\n")
            if len(value_lines) > 1:
                result_lines.append(f"  - {header}:")
                for vl in value_lines:
                    vl = vl.strip()
                    if vl:
                        result_lines.append(f"    - {vl}")
            else:
                result_lines.append(f"  - {header}: {value}")

    return "\n".join(result_lines) + "\n"


def _normalize_double_pipes(markdown: str) -> str:
    """Normalize || (double pipes) in table rows — pymupdf4llm artifact.

    Only handles the empty-cell case: |Cell1||Cell3| → |Cell1| |Cell3|
    For joined rows (separator||data), we try to split them into separate lines,
    but only when both halves look like complete pipe-table rows.
    """
    lines = markdown.split("\n")
    result = []
    for line in lines:
        stripped = line.strip()
        if "||" not in stripped or not stripped.startswith("|"):
            result.append(line)
            continue

        # Try to detect joined rows: |---|---||Data|Data|
        # Each half must start and end with | and look like a full row
        parts = stripped.split("||")
        if len(parts) == 2:
            left, right = parts
            left = left.strip()
            right = right.strip()
            # Ensure both halves are well-formed pipe rows
            if (left.startswith("|") and left.endswith("|") and
                    right and not right.startswith("|")):
                right = "|" + right
                if not right.endswith("|"):
                    right = right + "|"
                # Only split if one half is a separator
                if _is_separator(left) or _is_separator(right):
                    result.extend([left, right])
                    continue

        # Default: treat || as empty cells
        result.append(re.sub(r"\|\|", "| |", line))
    return "\n".join(result)


def _tables_to_lists(markdown: str) -> tuple:
    """Find markdown pipe-tables and convert them to nested bullet lists.

    Handles:
    - Standard tables (header + separator + data)
    - Tables with <br> in cells (from PDF extraction)
    - Tables split across pages (continuation with new separator)
    - Headerless/generic-header tables
    - Double-pipe || artifacts (empty cells)
    - CRLF line endings
    - Indented tables (leading whitespace before |)

    Returns (transformed_markdown, tables_found_count).
    """
    # Normalize line endings to LF
    markdown = markdown.replace("\r\n", "\n").replace("\r", "\n")

    # Pre-process: normalize || to | | before table detection
    markdown = _normalize_double_pipes(markdown)

    # Pattern: any block of consecutive pipe-rows (possibly with separator lines)
    # A table is 2+ consecutive lines starting with | (with optional leading whitespace)
    table_pattern = re.compile(
        r"((?:^[ \t]*\|.+\|[ \t]*$\n?){2,})",
        re.MULTILINE,
    )

    tables_found = 0

    def _replace(match):
        nonlocal tables_found
        block = match.group(1).strip()
        lines = [l for l in block.split("\n") if l.strip()]

        # Must have at least one non-separator line
        data_lines = [l for l in lines if not _is_separator(l)]
        if not data_lines:
            return block

        tables_found += 1
        return _table_to_list(lines)

    transformed = table_pattern.sub(_replace, markdown)

    # Attach table titles to their content: collapse blank lines between
    # a line that looks like a table caption and the bullet-list table body.
    # This prevents chunking from splitting the title from the table data.
    transformed = re.sub(
        r"(_[^_\n]+_)\s*\n\n+(- \*\*)",
        r"\1\n\2",
        transformed,
    )

    return transformed, tables_found


def _describe_image_with_vision(data_uri: str, llm_client) -> str:
    """Send a base64 image to Claude Vision and get a text description."""
    try:
        response = llm_client.chat.completions.create(
            messages=[{
                "role": "user",
                "content": [
                    {"type": "image_url", "image_url": {"url": data_uri}},
                    {"type": "text", "text": (
                        "Describe this image concisely for a technical document. "
                        "If it's a chart or diagram, describe the data and key takeaways. "
                        "If it's a UI screenshot, describe the interface elements and layout. "
                        "Keep it under 3 sentences."
                    )},
                ],
            }],
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        return f"Image could not be described: {e}"


def _handle_images(markdown: str, llm_client=None) -> tuple:
    """Replace image references with AI descriptions (if client available) or placeholders."""
    img_pattern = re.compile(r"!\[([^\]]*)\]\(([^)]+)\)")

    images_found = 0
    images_described = 0

    def _replace_img(match):
        nonlocal images_found, images_described
        images_found += 1
        alt_text = match.group(1).strip()
        url = match.group(2).strip()

        # If we have an LLM client and a base64 data URI, use Claude Vision
        if llm_client and url.startswith("data:image/"):
            images_described += 1
            desc = _describe_image_with_vision(url, llm_client)
            return f"[Image: {desc}]"

        # Fallback to alt text or generic placeholder
        if alt_text:
            return f"[Image: {alt_text}]"
        return "[Image: visual content]"

    transformed = img_pattern.sub(_replace_img, markdown)

    # pymupdf4llm placeholders: **==> picture [WxH] intentionally omitted <==**
    pymupdf_pattern = re.compile(r"\*?\*?=+>.*?picture\s*\[.*?\].*?omitted.*?<==\*?\*?")

    def _replace_pymupdf(match):
        nonlocal images_found
        images_found += 1
        return "[Image: visual content]"

    transformed = pymupdf_pattern.sub(_replace_pymupdf, transformed)
    return transformed, images_found


def process_document(file_path: str, llm_client=None) -> dict:
    """Full conversion pipeline: extract → tables to lists → image handling.

    Returns:
        {
            "raw_markdown": str,
            "tables_converted": str,
            "final_markdown": str,
            "stats": {"tables_found": int, "images_found": int, "char_count": int}
        }
    """
    # Validate file path is within temp directory to prevent path traversal
    real_path = os.path.realpath(file_path)
    temp_dir = os.path.realpath(tempfile.gettempdir())
    if not real_path.startswith(temp_dir) and not os.path.isabs(file_path):
        raise ValueError(f"Invalid file path: must be an absolute path or within temp directory")
    if not os.path.isfile(real_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    raw_markdown = _extract_markdown(real_path, llm_client=llm_client)

    tables_converted, tables_found = _tables_to_lists(raw_markdown)

    final_markdown, images_found = _handle_images(tables_converted, llm_client=llm_client)
    # Strip stray DOCX tab/sheet names (e.g. "Thẻ 1", "Thẻ 2")
    final_markdown = re.sub(r"^\s*Thẻ\s+\d+\s*$", "", final_markdown, flags=re.MULTILINE)
    # Collapse 3+ consecutive blank lines to 2
    final_markdown = re.sub(r"\n{4,}", "\n\n\n", final_markdown)

    return {
        "raw_markdown": raw_markdown,
        "tables_converted": tables_converted,
        "final_markdown": final_markdown,
        "stats": {
            "tables_found": tables_found,
            "images_found": images_found,
            "char_count": len(final_markdown),
        },
    }
